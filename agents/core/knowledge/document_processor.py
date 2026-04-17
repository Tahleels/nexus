"""
document_processor.py — Universal document ingestion pipeline.

Supported formats
-----------------
  .pdf         — pdfplumber   (pip install pdfplumber)
  .docx/.doc   — python-docx  (pip install python-docx)
  .xlsx/.xls   — openpyxl/pandas (already in requirements)
  .csv         — pandas       (already in requirements)
  .pptx/.ppt   — python-pptx  (already in requirements)
  .html/.htm   — beautifulsoup4 (pip install beautifulsoup4)
  .eml         — stdlib email
  .msg         — extract-msg  (pip install extract-msg)
  .png/.jpg... — easyocr      (pip install easyocr) [optional, CPU-only]
  .txt/.md     — stdlib
  .json        — stdlib

Pipeline
--------
  file → extract text + metadata
       → chunk (~500 tokens, 50-token overlap, respects paragraphs)
       → embed + store chunks in ChromaDB (documents collection)
       → store metadata in SQL Server app_documents table
       → return {document_id, filename, file_type, chunk_count, ...}

Azure migration
---------------
  Extract step : swap per-extractor calls with Azure AI Document Intelligence
                 (Form Recognizer) — returns the same (text, meta) tuple.
  Store step   : unchanged — chromadb / vector_store.py stays the same.
  SQL meta     : unchanged — becomes Azure SQL Database connection.
"""
import hashlib
import json
from logging_config import get_logger
import os
import re
import sys
import uuid
from typing import Dict, List, Optional, Tuple

logger = get_logger(__name__)

# ── Project root ──────────────────────────────────────────────────────────────
_HERE         = os.path.dirname(os.path.abspath(__file__))
_PROJECT_ROOT = os.path.normpath(os.path.join(_HERE, '..', '..', '..'))
if _PROJECT_ROOT not in sys.path:
    sys.path.insert(0, _PROJECT_ROOT)

from llm_providers.factory import get_provider
from llm_providers.errors import LLMError
from llm_providers.models import resolve_default_model

_env_store = os.environ.get('AGENT_FILE_STORE', '').strip()
_STORE_DIR  = _env_store if _env_store else os.path.join(_PROJECT_ROOT, 'Data', 'agent_store')
_RAW_DIR    = os.path.join(_STORE_DIR, 'raw_documents')
os.makedirs(_RAW_DIR, exist_ok=True)

# Chunking constants (~500 tokens at 4 chars/token, 50-token overlap)
_CHUNK_CHARS   = 2000
_OVERLAP_CHARS = 200

# LLM models used for metric extraction / summarization (see llm_providers)
_METRICS_MODEL = resolve_default_model("DOCUMENT_METRICS_MODEL", tier="fast")
_SUMMARY_MODEL = resolve_default_model("DOCUMENT_SUMMARY_MODEL", tier="fast")


# ═══════════════════════════════════════════════════════════════════════════════
# CHUNKER
# ═══════════════════════════════════════════════════════════════════════════════

def _chunk_text(text: str) -> List[str]:
    """Split text into overlapping chunks, preferring paragraph/sentence breaks."""
    text = text.strip()
    if not text:
        return []
    if len(text) <= _CHUNK_CHARS:
        return [text]

    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = start + _CHUNK_CHARS
        if end >= len(text):
            chunk = text[start:].strip()
            if chunk:
                chunks.append(chunk)
            break
        # Prefer paragraph > sentence > word boundary
        for sep in ('\n\n', '. ', ' '):
            bp = text.rfind(sep, start + _OVERLAP_CHARS, end)
            if bp != -1 and bp > start:
                end = bp + len(sep)
                break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = max(end - _OVERLAP_CHARS, start + 1)

    return [c for c in chunks if c]


# ═══════════════════════════════════════════════════════════════════════════════
# EXTRACTORS  — each returns (text: str, meta: dict)
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_pdf(path: str) -> Tuple[str, Dict]:
    """
    Extract text from a PDF via pdfplumber, page by page.

    Args:
        path: Absolute path to the .pdf file.

    Returns:
        (text, meta) where text is all pages joined with blank lines and
        meta contains {"page_count": int}, or ("", {"error": ...}) on failure
        (including pdfplumber not being installed).
    """
    try:
        import pdfplumber
        pages: List[str] = []
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
        return "\n\n".join(pages), {"page_count": page_count}
    except ImportError:
        logger.warning("pdfplumber not installed — run: pip install pdfplumber")
        return "", {"error": "pdfplumber not installed"}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_word(path: str) -> Tuple[str, Dict]:
    """
    Extract text from a Word document, dispatching by extension.

    .doc files (legacy binary OLE format) cannot be read by python-docx, so
    they are routed to `_extract_doc_legacy`. .docx files go straight to
    `_parse_docx_file`.

    Args:
        path: Absolute path to the .doc or .docx file.

    Returns:
        (text, meta) tuple — see `_parse_docx_file` / `_extract_doc_legacy`.
    """
    ext = os.path.splitext(path)[1].lower()

    # .doc (old binary OLE format) — python-docx cannot read it
    if ext == ".doc":
        return _extract_doc_legacy(path)

    return _parse_docx_file(path)


def _parse_docx_file(path: str) -> Tuple[str, Dict]:
    """
    Parse a .docx file with python-docx, extracting paragraph and table text.

    Args:
        path: Absolute path to a .docx file (or a .docx produced by
            converting a legacy .doc via LibreOffice).

    Returns:
        (text, meta) where text is paragraphs followed by table rows
        (cells joined with " | "), separated by blank lines, and meta
        contains {"paragraph_count": int}. Returns ("", {"error": ...}) if
        python-docx is missing or parsing fails.
    """
    try:
        from docx import Document
        doc   = Document(path)
        parts: List[str] = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts), {"paragraph_count": len(parts)}
    except ImportError:
        logger.warning("python-docx not installed — run: pip install python-docx")
        return "", {"error": "python-docx not installed"}
    except Exception as exc:
        return "", {"error": str(exc)}


def _find_soffice() -> Optional[str]:
    """
    Locate the headless LibreOffice binary used to convert legacy .doc files.

    Returns:
        Absolute path to soffice/soffice.exe if found on PATH or in one of the
        two hardcoded Program Files locations, else None.
    """
    import shutil
    found = shutil.which("soffice") or shutil.which("soffice.exe")
    if found:
        return found
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if os.path.isfile(candidate):
            return candidate
    return None


def _extract_doc_legacy(path: str) -> Tuple[str, Dict]:
    """
    Extract text from old binary .doc files. Tries Word COM automation first
    (only works if Word is installed — rare on a server), then falls back to
    headless LibreOffice conversion to .docx, then docx2txt as a last resort.
    """
    import tempfile, shutil, subprocess

    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        abs_path = os.path.abspath(path)
        try:
            doc = word.Documents.Open(abs_path, ReadOnly=True)
            text = doc.Content.Text
            doc.Close(False)
        finally:
            word.Quit()
        text = text.strip()
        if text:
            return text, {"source": "win32com", "paragraph_count": text.count("\r")}
    except ImportError:
        pass
    except Exception as exc:
        # -2147221005 == CO_E_CLASSSTRING: Word not installed/registered — expected on servers
        args = getattr(exc, "args", ())
        if args and args[0] == -2147221005:
            logger.debug("win32com: Word not registered on this machine — skipping to LibreOffice")
        else:
            logger.warning("win32com .doc extraction failed: %s — trying LibreOffice", exc)

    # Fallback: headless LibreOffice conversion to .docx, then parse normally
    soffice = _find_soffice()
    if soffice:
        out_dir = tempfile.mkdtemp(prefix="doc2docx_")
        try:
            subprocess.run(
                [soffice, "--headless", "--norestore", "--convert-to", "docx",
                 "--outdir", out_dir, os.path.abspath(path)],
                capture_output=True, timeout=120, check=False,
            )
            base      = os.path.splitext(os.path.basename(path))[0]
            converted = os.path.join(out_dir, f"{base}.docx")
            if os.path.isfile(converted):
                text, meta = _parse_docx_file(converted)
                if text.strip():
                    meta["source"] = "libreoffice"
                    return text, meta
        except Exception as exc:
            logger.warning("LibreOffice .doc conversion failed: %s", exc)
        finally:
            shutil.rmtree(out_dir, ignore_errors=True)
    else:
        logger.warning(".doc extraction: soffice (LibreOffice) not found on PATH")

    # Last resort: docx2txt (handles some .doc files actually saved as OOXML)
    try:
        import docx2txt
        text = docx2txt.process(path) or ""
        if text.strip():
            return text.strip(), {"source": "docx2txt"}
    except Exception:
        pass

    return "", {"error": ".doc extraction failed — LibreOffice not found and Word/win32com unavailable"}


def _extract_excel(path: str) -> Tuple[str, Dict]:
    """
    Extract text from an Excel workbook (.xlsx/.xls/.xlsm) via pandas.

    Reads up to the first 10 sheets (up to 500 rows each), rendering each
    sheet as a "Sheet: <name>" header followed by comma-joined rows.

    Args:
        path: Absolute path to the Excel file.

    Returns:
        (text, meta) where meta contains {"sheet_count": int, "sheets": [...]}
        or ("", {"error": ...}) if every sheet is empty or reading fails.
    """
    try:
        import pandas as pd
        xl          = pd.ExcelFile(path)
        sheet_names = xl.sheet_names
        parts: List[str] = []
        for sheet in sheet_names[:10]:
            df = pd.read_excel(xl, sheet_name=sheet, nrows=500)
            if df.empty:
                continue
            lines = [", ".join(str(c) for c in df.columns)]
            for _, row in df.iterrows():
                lines.append(", ".join(str(v) for v in row.values))
            parts.append(f"Sheet: {sheet}\n" + "\n".join(lines))
        if not parts:
            return "", {"error": "All sheets are empty"}
        return "\n\n".join(parts), {"sheet_count": len(sheet_names), "sheets": sheet_names[:10]}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_csv(path: str) -> Tuple[str, Dict]:
    """
    Extract text from a CSV file, trying multiple encodings in order.

    Reads up to 1000 rows and renders them as comma-joined lines (header
    row first). Bad lines are skipped rather than raising.

    Args:
        path: Absolute path to the .csv file.

    Returns:
        (text, meta) where meta contains {"row_count", "columns", "encoding"},
        or ("", {"error": ...}) if no encoding in the fallback chain succeeds.
    """
    import pandas as pd
    # Try encodings in order: utf-8-sig handles Excel BOM, latin-1 handles Windows cp1252
    for enc in ("utf-8-sig", "utf-8", "latin-1", "cp1252"):
        try:
            df = pd.read_csv(path, nrows=1000, encoding=enc, on_bad_lines="skip")
            if df.empty:
                return "", {"error": "CSV parsed but is empty", "encoding": enc}
            # Build readable text: header + rows
            lines = [", ".join(str(c) for c in df.columns)]
            for _, row in df.iterrows():
                lines.append(", ".join(str(v) for v in row.values))
            text = "\n".join(lines)
            return text, {"row_count": len(df), "columns": list(df.columns), "encoding": enc}
        except Exception as exc:
            last_exc = exc
            continue
    return "", {"error": str(last_exc)}


def _extract_ppt(path: str) -> Tuple[str, Dict]:
    """
    Extract text from a PowerPoint deck (.ppt/.pptx) via python-pptx.

    Args:
        path: Absolute path to the presentation file.

    Returns:
        (text, meta) where text lists each slide's shape text under a
        "Slide N:" header, and meta contains {"slide_count": int}.
        Returns ("", {"error": ...}) on failure.
    """
    try:
        from pptx import Presentation
        prs   = Presentation(path)
        parts: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                s.text.strip()
                for s in slide.shapes
                if hasattr(s, "text") and s.text.strip()
            ]
            if texts:
                parts.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(parts), {"slide_count": len(prs.slides)}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_html(path: str) -> Tuple[str, Dict]:
    """
    Extract visible text from an HTML/HTM file via BeautifulSoup.

    Strips script/style/nav/footer/header tags before extracting text, then
    collapses runs of 3+ blank lines down to a single blank line.

    Args:
        path: Absolute path to the .html/.htm file.

    Returns:
        (text, {"source": "html"}), or ("", {"error": ...}) if
        beautifulsoup4 is missing or parsing fails.
    """
    try:
        from bs4 import BeautifulSoup
        with open(path, "r", encoding="utf-8", errors="replace") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = re.sub(r"\n{3,}", "\n\n", soup.get_text(separator="\n")).strip()
        return text, {"source": "html"}
    except ImportError:
        logger.warning("beautifulsoup4 not installed — run: pip install beautifulsoup4")
        return "", {"error": "beautifulsoup4 not installed"}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_eml(path: str) -> Tuple[str, Dict]:
    """
    Extract text from a .eml email file using the stdlib `email` module.

    Builds a "From / Date / Subject" header followed by the first
    text/plain body part found (or the whole payload for non-multipart
    messages).

    Args:
        path: Absolute path to the .eml file.

    Returns:
        (text, meta) where meta contains {"subject", "sender", "date"},
        or ("", {"error": ...}) on failure.
    """
    import email as _email_mod
    try:
        with open(path, "rb") as fh:
            msg = _email_mod.message_from_bytes(fh.read())
        subject = msg.get("Subject", "")
        sender  = msg.get("From", "")
        date    = msg.get("Date", "")
        body    = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_payload(decode=True).decode("utf-8", errors="replace")
                    break
        else:
            body = msg.get_payload(decode=True).decode("utf-8", errors="replace")
        text = f"From: {sender}\nDate: {date}\nSubject: {subject}\n\n{body}"
        return text, {"subject": subject, "sender": sender, "date": date}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_msg(path: str) -> Tuple[str, Dict]:
    """
    Extract text from an Outlook .msg file via extract-msg.

    Args:
        path: Absolute path to the .msg file.

    Returns:
        (text, meta) where meta contains {"subject", "sender", "date"},
        or ("", {"error": ...}) if extract-msg is missing or parsing fails.
    """
    try:
        import extract_msg
        msg     = extract_msg.Message(path)
        subject = msg.subject or ""
        sender  = msg.sender  or ""
        date    = str(msg.date or "")
        body    = msg.body    or ""
        text    = f"From: {sender}\nDate: {date}\nSubject: {subject}\n\n{body}"
        return text, {"subject": subject, "sender": sender, "date": date}
    except ImportError:
        logger.warning("extract-msg not installed — run: pip install extract-msg")
        return "", {"error": "extract-msg not installed"}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_image(path: str) -> Tuple[str, Dict]:
    """
    OCR an image file (.png/.jpg/.jpeg/.tiff/.bmp) via easyocr (CPU-only).

    Args:
        path: Absolute path to the image file.

    Returns:
        (text, meta) where text is the OCR'd lines joined with newlines and
        meta contains {"source": "easyocr", "size": "<w>x<h>"}, or
        ("", {"error": ...}) if easyocr is missing or OCR fails.
    """
    try:
        import easyocr
        from PIL import Image
        img    = Image.open(path)
        # gpu=False keeps it CPU-only so it works on any server without CUDA
        reader = easyocr.Reader(['en'], gpu=False, verbose=False)
        result = reader.readtext(path, detail=0)
        text   = "\n".join(result)
        return text, {"source": "easyocr", "size": f"{img.width}x{img.height}"}
    except ImportError:
        logger.warning("easyocr not installed — run: pip install easyocr")
        return "", {"error": "easyocr not installed"}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_text(path: str) -> Tuple[str, Dict]:
    """
    Read a .txt/.md file as plain UTF-8 text (invalid bytes replaced).

    Args:
        path: Absolute path to the text file.

    Returns:
        (text, {}) on success, or ("", {"error": ...}) on failure.
    """
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as fh:
            return fh.read(), {}
    except Exception as exc:
        return "", {"error": str(exc)}


def _extract_json(path: str) -> Tuple[str, Dict]:
    """
    Read a .json file and re-serialize it with indent=2 for readability.

    Args:
        path: Absolute path to the .json file.

    Returns:
        (text, {"source": "json"}) on success, or ("", {"error": ...}) if
        the file is missing or not valid JSON.
    """
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as fh:
            data = json.load(fh)
        return json.dumps(data, indent=2, ensure_ascii=False), {"source": "json"}
    except Exception as exc:
        return "", {"error": str(exc)}


_EXTRACTOR_MAP = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_word,
    ".doc":  _extract_word,
    ".xlsx": _extract_excel,
    ".xls":  _extract_excel,
    ".xlsm": _extract_excel,
    ".csv":  _extract_csv,
    ".pptx": _extract_ppt,
    ".ppt":  _extract_ppt,
    ".html": _extract_html,
    ".htm":  _extract_html,
    ".eml":  _extract_eml,
    ".msg":  _extract_msg,
    ".png":  _extract_image,
    ".jpg":  _extract_image,
    ".jpeg": _extract_image,
    ".tiff": _extract_image,
    ".bmp":  _extract_image,
    ".txt":  _extract_text,
    ".md":   _extract_text,
    ".json": _extract_json,
}

SUPPORTED_EXTENSIONS = sorted(_EXTRACTOR_MAP.keys())


def extract_text(path: str) -> Tuple[str, Dict]:
    """Extract plain text from a file. Returns (text, metadata) for any supported extension."""
    ext = os.path.splitext(path)[1].lower()
    extractor = _EXTRACTOR_MAP.get(ext)
    if not extractor:
        return "", {"error": f"Unsupported extension: {ext}"}
    return extractor(path)


# ═══════════════════════════════════════════════════════════════════════════════
# KEY METRICS EXTRACTOR
# ═══════════════════════════════════════════════════════════════════════════════

def _select_metric_chunks(chunks: List[str], key_metrics: List[str],
                          max_chars: int = 4000) -> str:
    """
    Score each chunk by how many metric keywords it contains, then pick the
    top-scored chunks that fit within max_chars.  Falls back to the first two
    chunks when nothing scores above zero.
    """
    keywords = [m.lower() for m in key_metrics]

    def score(chunk: str) -> int:
        lower = chunk.lower()
        return sum(1 for kw in keywords if kw in lower)

    scored = sorted(enumerate(chunks), key=lambda x: score(x[1]), reverse=True)

    selected, total = [], 0
    for _, chunk in scored:
        if total >= max_chars:
            break
        take = chunk[:max_chars - total]
        selected.append(take)
        total += len(take)

    # If no chunk scored anything, fall back to the first two chunks
    if not any(score(c) > 0 for c in chunks):
        selected = []
        total = 0
        for chunk in chunks[:2]:
            take = chunk[:max_chars - total]
            selected.append(take)
            total += len(take)
            if total >= max_chars:
                break

    return "\n\n---\n\n".join(selected)


def _extract_key_metrics(chunks: List[str], key_metrics: List[str]) -> Dict:
    """
    Select the most relevant chunks via keyword scoring, then ask the LLM to
    extract the named metrics from that small context only.
    Returns a dict of {metric: value_or_null}.  Never raises.
    """
    if not key_metrics or not chunks:
        return {}

    context      = _select_metric_chunks(chunks, key_metrics)
    metrics_list = ", ".join(key_metrics)
    prompt = (
        f"Extract the following key metrics from the document excerpts below: {metrics_list}\n\n"
        f"Return ONLY a valid JSON object where keys are the metric names and values are "
        f"the extracted information as strings. Use null for any metric not found.\n\n"
        f"Document excerpts:\n{context}"
    )
    try:
        result = get_provider().chat(
            system="You are a document analysis assistant. Extract specific fields and return valid JSON only.",
            messages=[{"role": "user", "content": prompt}],
            model=_METRICS_MODEL,
            temperature=0,
            response_format={"type": "json_object"},
        )
        return json.loads(result.content)
    except LLMError as exc:
        logger.warning("_extract_key_metrics failed: %s", exc)
        return {}
    except Exception as exc:
        logger.warning("_extract_key_metrics failed: %s", exc)
        return {}


# ═══════════════════════════════════════════════════════════════════════════════
# DOCUMENT PROFILE HELPERS  (summary + profile table for Stage 1 retrieval)
# ═══════════════════════════════════════════════════════════════════════════════

def _generate_summary(text: str) -> str:
    """Generate a 2-3 sentence summary via GPT-4o-mini. Falls back to first 300 chars."""
    if not text:
        return text[:300].strip()
    context = text[:3000]
    prompt = (
        "Write a 2-3 sentence summary of this document. "
        "Include: document type, main subject/person/topic, and the key skills, "
        "qualifications, dates, or content areas covered. Be specific and factual.\n\n"
        f"Document:\n{context}"
    )
    try:
        result = get_provider().chat(
            system="You are a document summarization assistant. Return only the summary, no preamble.",
            messages=[{"role": "user", "content": prompt}],
            model=_SUMMARY_MODEL,
            temperature=0,
            max_tokens=150,
        )
        return (result.content or "").strip()
    except LLMError as exc:
        logger.warning("_generate_summary failed: %s", exc)
        return text[:300].strip()
    except Exception as exc:
        logger.warning("_generate_summary failed: %s", exc)
        return text[:300].strip()


def _embed_summary(text: str) -> Optional[list]:
    """Embed arbitrary text for Stage 1 similarity ranking. Returns None on failure."""
    if not text:
        return None
    try:
        from agents.core.knowledge.vector_store import _embed_one
        return _embed_one(text)
    except Exception as exc:
        logger.warning("_embed_summary failed: %s", exc)
        return None


def _build_embedding_text(summary: str, key_metrics: Optional[Dict]) -> str:
    """
    Combine the prose summary with ALL extracted key_metrics values before
    embedding. The narrative summary alone can under-represent specific
    structured details (e.g. specialized skills on a resume, a renewal
    clause on a contract, a spec on a product sheet) relative to a
    document's broader framing — appending every extracted metric value
    ensures specific facts directly influence the embedding instead of
    depending on the summary having mentioned them. This is domain-agnostic:
    it uses whatever metric fields were configured for extraction on that
    connector, so it works the same way for resumes, contracts, product
    catalogs, or any other document type — not resume-specific.
    """
    parts = [summary or ""]
    if key_metrics:
        metric_text = "; ".join(f"{k}: {v}" for k, v in key_metrics.items() if v)
        if metric_text:
            parts.append(metric_text)
    return "\n".join(p for p in parts if p)


def _upsert_document_profile(
    doc_id: str,
    source_file: str,
    summary: str,
    key_metrics: Optional[Dict],
    source_type: str = "knowledge",
    connector_key: Optional[str] = None,
    document_type: str = "general",
) -> None:
    """Upsert one row per document into document_profiles. Never raises."""
    try:
        from app_db import get_app_db
        embedding_text = _build_embedding_text(summary, key_metrics)
        embedding = _embed_summary(embedding_text)
        emb_json  = json.dumps(embedding) if embedding else None
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("""
                MERGE document_profiles AS t
                USING (SELECT ? AS doc_id) AS s ON t.doc_id = s.doc_id
                WHEN MATCHED THEN UPDATE SET
                    summary = ?, key_metrics = ?, summary_embedding = ?,
                    processed_at = GETUTCDATE()
                WHEN NOT MATCHED THEN INSERT
                    (doc_id, source_file, summary, key_metrics, summary_embedding,
                     source_type, connector_key, document_type)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?);
            """,
            doc_id,
            summary, json.dumps(key_metrics) if key_metrics else None, emb_json,
            doc_id, source_file, summary,
            json.dumps(key_metrics) if key_metrics else None, emb_json,
            source_type, connector_key, document_type)
            conn.commit()
    except Exception as exc:
        logger.warning("_upsert_document_profile failed: %s", exc)


# ═══════════════════════════════════════════════════════════════════════════════
# SQL METADATA HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _find_existing_doc_id(file_hash: str, uploaded_by_id: int,
                          scope: str, scope_id: Optional[int]) -> Optional[str]:
    """
    Return the existing doc_id for this exact (content, owner, scope) combination,
    or None if not yet ingested.

    Matching on owner + scope (not just content hash) keeps documents isolated:
      - different users uploading identical content always get separate docs
        (separate visibility — never merge across uploaded_by_id)
      - a user-scoped copy and a department-scoped copy of the same bytes
        never collapse into one document, even for the same uploader
    Including scope/scope_id also covers system/connector uploads
    (uploaded_by_id=0): re-scanning the same department- or global-scoped file
    refreshes it in place instead of creating a new duplicate every scan.
    """
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            if scope_id is None:
                cur.execute(
                    "SELECT TOP 1 id FROM app_documents "
                    "WHERE file_hash = ? AND uploaded_by_id = ? AND scope = ? "
                    "AND scope_id IS NULL ORDER BY created_at DESC",
                    file_hash, uploaded_by_id, scope,
                )
            else:
                cur.execute(
                    "SELECT TOP 1 id FROM app_documents "
                    "WHERE file_hash = ? AND uploaded_by_id = ? AND scope = ? AND scope_id = ? "
                    "ORDER BY created_at DESC",
                    file_hash, uploaded_by_id, scope, scope_id,
                )
            row = cur.fetchone()
        return row[0] if row else None
    except Exception as exc:
        logger.warning("_find_existing_doc_id: %s", exc)
        return None


def _save_doc_meta(doc_id: str, filename: str, file_type: str,
                   client_id: str, source_name: str, chunk_count: int,
                   file_hash: str, extra_meta: Dict,
                   uploaded_by_id: int = 0, uploaded_by_name: str = "",
                   version: int = 1, parent_doc_id: Optional[str] = None,
                   status: str = "ready",
                   scope: str = "user", scope_id: Optional[int] = None,
                   source_watch_id: Optional[int] = None,
                   source_watch_type: str = "") -> None:
    """
    Upsert document metadata into the SQL Server `app_documents` table.

    Uses a MERGE: inserts a new row keyed by `doc_id` if absent, otherwise
    updates only chunk_count/status/extra_meta/scope/scope_id (the columns
    that legitimately change on re-ingestion of the same doc_id). Never
    raises — failures are logged and swallowed so a metadata-write hiccup
    doesn't block ingestion.
    """
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("""
                MERGE app_documents AS t
                USING (SELECT ? AS doc_id) AS s ON t.id = s.doc_id
                WHEN MATCHED THEN UPDATE SET
                    chunk_count = ?, status = ?, extra_meta = ?,
                    scope = ?, scope_id = ?
                WHEN NOT MATCHED THEN
                    INSERT (id, filename, file_type, client_id, source_name,
                            chunk_count, file_hash, extra_meta,
                            uploaded_by_id, uploaded_by_name,
                            version, parent_doc_id, status, scope, scope_id,
                            source_watch_id, source_watch_type, created_at)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?,
                            ?, ?, GETUTCDATE());
            """, doc_id,
                chunk_count, status, json.dumps(extra_meta), scope, scope_id,
                doc_id, filename, file_type, client_id or "",
                source_name or filename, chunk_count, file_hash,
                json.dumps(extra_meta), uploaded_by_id, uploaded_by_name,
                version, parent_doc_id, status, scope, scope_id,
                source_watch_id, source_watch_type or None)
            conn.commit()
    except Exception as exc:
        logger.warning("document_processor: metadata save failed: %s", exc)


def _delete_doc_meta(doc_id: str) -> None:
    """Delete a document's row from `app_documents`. Logs and swallows errors."""
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("DELETE FROM app_documents WHERE id = ?", doc_id)
            conn.commit()
    except Exception as exc:
        logger.warning("document_processor: metadata delete failed: %s", exc)


# ═══════════════════════════════════════════════════════════════════════════════
# PUBLIC API
# ═══════════════════════════════════════════════════════════════════════════════

def process_document(file_path: str,
                     client_id:         str = "",
                     source_name:       str = "",
                     uploaded_by_id:    int = 0,
                     uploaded_by_name:  str = "",
                     version:           int = 1,
                     parent_doc_id:     Optional[str] = None,
                     scope:             str = "user",
                     scope_id:          Optional[int] = None,
                     source_watch_id:   Optional[int] = None,
                     source_watch_type: str = "",
                     key_metrics:       Optional[List[str]] = None) -> Dict:
    """
    Ingest a document into the vector knowledge base.

    Parameters
    ----------
    file_path        : Absolute path, or filename relative to raw_documents store.
    client_id        : Optional tenant/client label for scoping searches.
    source_name      : Human label (e.g. "Q1 Sales Report").
    uploaded_by_id   : User ID of the uploader (0 = system/anonymous).
    uploaded_by_name : Display name of the uploader.
    version          : Version number (1 = first upload, 2+ = replacement).
    parent_doc_id    : Previous document ID when uploading a new version.
    scope            : 'user' | 'global' | 'department' | 'project'
    scope_id         : dept_id or project_id when scope='department'/'project'

    Returns
    -------
    {success, document_id, filename, file_type, chunk_count, text_length,
     vector_store_available, extraction_warning?}
    """
    # ── Resolve path ──────────────────────────────────────────────────────────
    if not os.path.isabs(file_path):
        candidate = os.path.join(_RAW_DIR, file_path)
        if os.path.exists(candidate):
            file_path = candidate

    if not os.path.exists(file_path):
        return {"success": False, "error": f"File not found: {file_path}"}

    filename  = os.path.basename(file_path)
    ext       = os.path.splitext(filename)[1].lower()
    extractor = _EXTRACTOR_MAP.get(ext)
    if not extractor:
        return {
            "success": False,
            "error":   f"Unsupported file type '{ext}'.",
            "supported_extensions": SUPPORTED_EXTENSIONS,
        }

    # ── File hash (dedup key) ─────────────────────────────────────────────────
    try:
        with open(file_path, "rb") as fh:
            file_hash = hashlib.sha256(fh.read()).hexdigest()[:16]
    except Exception:
        file_hash = uuid.uuid4().hex[:16]

    # ── Dedup: reuse existing doc_id for the same (content, owner, scope) ──────
    # Same file from chat/SharePoint/directory for the same user+scope → one doc_id.
    # Different users, or different scope, uploading the same file → separate documents.
    _scope = scope or "user"
    doc_id = (_find_existing_doc_id(file_hash, uploaded_by_id, _scope, scope_id)
              or f"doc_{file_hash}_{uuid.uuid4().hex[:8]}")

    # ── Extract ───────────────────────────────────────────────────────────────
    text, extra_meta = extractor(file_path)
    if not text or not text.strip():
        return {
            "success":          False,
            "document_id":      doc_id,
            "filename":         filename,
            "error":            "No text could be extracted.",
            "extraction_meta":  extra_meta,
        }

    # ── Chunk ─────────────────────────────────────────────────────────────────
    raw_chunks = _chunk_text(text)

    # ── Extract key metrics from the most relevant chunks only ────────────────
    metrics_result: Dict = {}
    if key_metrics:
        metrics_result = _extract_key_metrics(raw_chunks, key_metrics)
        if metrics_result:
            extra_meta["extracted_metrics"] = metrics_result

    # ── Generate document summary for Stage 1 comprehensive retrieval ─────────
    _summary = _generate_summary(text)

    chunks = [
        {
            "text": chunk,
            "metadata": {
                "document_id": doc_id,
                "filename":    filename,
                "file_type":   ext.lstrip("."),
                "client_id":   client_id or "",
                "source_name": source_name or filename,
                "chunk_index": str(i),
                "chunk_total": str(len(raw_chunks)),
            },
        }
        for i, chunk in enumerate(raw_chunks)
    ]

    # ── Store in ChromaDB ─────────────────────────────────────────────────────
    from agents.core.knowledge.vector_store import get_vector_store
    vs      = get_vector_store()
    stored  = vs.add_chunks(doc_id, chunks)

    # Persist the raw file path so delete_document can clean it up from disk
    extra_meta["_raw_file"] = file_path

    # ── Store metadata in SQL Server ──────────────────────────────────────────
    _save_doc_meta(
        doc_id=doc_id, filename=filename, file_type=ext.lstrip("."),
        client_id=client_id, source_name=source_name or filename,
        chunk_count=stored, file_hash=file_hash, extra_meta=extra_meta,
        uploaded_by_id=uploaded_by_id, uploaded_by_name=uploaded_by_name,
        version=version, parent_doc_id=parent_doc_id,
        scope=scope or "user", scope_id=scope_id,
        source_watch_id=source_watch_id,
        source_watch_type=source_watch_type,
    )

    # ── Upsert document profile for comprehensive listing ─────────────────────
    _connector_key = (f"{source_watch_type}:{source_watch_id}"
                      if source_watch_id and source_watch_type else None)
    _upsert_document_profile(
        doc_id=doc_id,
        source_file=filename,
        summary=_summary,
        key_metrics=metrics_result or None,
        source_type="connector" if _connector_key else "knowledge",
        connector_key=_connector_key,
        document_type="general",
    )

    result: Dict = {
        "success":               True,
        "document_id":           doc_id,
        "filename":              filename,
        "file_type":             ext.lstrip("."),
        "chunk_count":           stored,
        "text_length":           len(text),
        "vector_store_available": vs.available,
    }
    if extra_meta.get("error"):
        result["extraction_warning"] = extra_meta["error"]
    logger.info("process_document: %s -> %d chunks (id=%s)", filename, stored, doc_id)
    return result


def search_documents(query: str, n_results: int = 5,
                     client_id: Optional[str] = None,
                     doc_ids: Optional[set] = None) -> Dict:
    """
    Semantic search over ingested document chunks.

    Parameters
    ----------
    query     : Natural-language question or keywords.
    n_results : Maximum number of chunks to return.
    client_id : Restrict to documents from a specific client/tenant.
    doc_ids   : Restrict to specific document IDs (SQL-level filter — always
                finds chunks from these docs regardless of global cosine ranking).
    """
    from agents.core.knowledge.vector_store import get_vector_store
    vs = get_vector_store()

    if not vs.available:
        return {
            "success": False,
            "query":   query,
            "count":   0,
            "results": [],
            "error":   "Vector store unavailable.",
        }

    where = {"client_id": client_id} if client_id else None
    hits  = vs.search_documents(query, n_results=n_results, where=where, doc_ids=doc_ids)

    return {
        "success": True,
        "query":   query,
        "count":   len(hits),
        "results": hits,
    }


def get_visible_doc_ids(user_id: int, role: str) -> Optional[List[str]]:
    """
    Return the set of doc IDs visible to this user, or None (= all) for admins/devs.

    Visibility rules:
      scope='global'     — visible to everyone
      scope='user'       — visible to uploader + explicitly assigned users
      scope='department' — visible to all members of that department
      scope='project'    — visible to all members of that project
      Direct assignments always grant access regardless of scope.
    """
    if role in ("admin", "dev"):
        return None  # no filter — see everything

    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()

            # Collect org memberships
            cur.execute("SELECT dept_id FROM user_departments WHERE user_id = ?", user_id)
            dept_ids    = [r[0] for r in cur.fetchall()]
            cur.execute("SELECT project_id FROM user_projects WHERE user_id = ?", user_id)
            project_ids = [r[0] for r in cur.fetchall()]

            # Base: global docs + own user-scoped docs + directly assigned docs
            cur.execute("""
                SELECT id FROM app_documents
                WHERE scope = 'global'
                   OR (scope = 'user' AND uploaded_by_id = ?)
                UNION
                SELECT doc_id FROM app_document_assignments WHERE user_id = ?
            """, user_id, user_id)
            ids = {r[0] for r in cur.fetchall()}

            # Dept-scoped docs
            if dept_ids:
                ph = ",".join("?" for _ in dept_ids)
                cur.execute(f"""
                    SELECT id FROM app_documents
                    WHERE scope = 'department' AND scope_id IN ({ph})
                """, *dept_ids)
                ids.update(r[0] for r in cur.fetchall())

            # Project-scoped docs
            if project_ids:
                ph = ",".join("?" for _ in project_ids)
                cur.execute(f"""
                    SELECT id FROM app_documents
                    WHERE scope = 'project' AND scope_id IN ({ph})
                """, *project_ids)
                ids.update(r[0] for r in cur.fetchall())

            return list(ids)
    except Exception as exc:
        logger.warning("get_visible_doc_ids failed: %s", exc)
        return []


def list_documents(user_id: int = 0, role: str = "user",
                   client_id: Optional[str] = None,
                   source_watch_type: Optional[str] = None,
                   source_watch_id: Optional[int] = None) -> List[Dict]:
    """
    List documents visible to this user.
    Admins/devs see all; regular users see only their own + assigned docs.
    Optionally filter by source_watch_type ('filesystem', 'sharepoint', 'upload')
    and/or source_watch_id.
    """
    try:
        from app_db import get_app_db
        visible_ids = get_visible_doc_ids(user_id, role)
        with get_app_db() as conn:
            cur = conn.cursor()
            select_cols = """
                    SELECT d.id, d.filename, d.file_type, d.source_name,
                           d.chunk_count, d.uploaded_by_id, d.uploaded_by_name,
                           d.version, d.parent_doc_id, d.status, d.created_at,
                           d.scope, d.scope_id, d.extra_meta,
                           d.source_watch_type, d.source_watch_id,
                           (SELECT STRING_AGG(u.user_name, ', ')
                            FROM app_document_assignments u WHERE u.doc_id = d.id
                           ) AS assigned_to
                    FROM app_documents d
            """
            # Build extra WHERE clauses for source filters
            src_clauses, src_params = [], []
            if source_watch_type == "upload":
                src_clauses.append("d.source_watch_type IS NULL")
            elif source_watch_type in ("filesystem", "sharepoint"):
                src_clauses.append("d.source_watch_type = ?")
                src_params.append(source_watch_type)
            if source_watch_id is not None:
                src_clauses.append("d.source_watch_id = ?")
                src_params.append(source_watch_id)

            def _src_and(existing_where: bool) -> str:
                if not src_clauses:
                    return ""
                kw = " AND " if existing_where else " WHERE "
                return kw + " AND ".join(src_clauses)

            if visible_ids is None:
                # admin: no restriction
                if client_id:
                    q = select_cols + " WHERE d.client_id = ?" + _src_and(True) + " ORDER BY d.created_at DESC"
                    cur.execute(q, client_id, *src_params)
                else:
                    q = select_cols + _src_and(False) + " ORDER BY d.created_at DESC"
                    cur.execute(q, *src_params)
            elif not visible_ids:
                return []
            else:
                placeholders = ",".join(["?" for _ in visible_ids])
                base = select_cols + f" WHERE d.id IN ({placeholders})"
                params = list(visible_ids)
                if client_id:
                    q = base + " AND d.client_id = ?" + _src_and(True) + " ORDER BY d.created_at DESC"
                    cur.execute(q, *params, client_id, *src_params)
                else:
                    q = base + _src_and(True) + " ORDER BY d.created_at DESC"
                    cur.execute(q, *params, *src_params)
            cols = [d[0] for d in cur.description]
            return [dict(zip(cols, row)) for row in cur.fetchall()]
    except Exception as exc:
        logger.warning("list_documents failed: %s", exc)
        return []


def get_document(doc_id: str, user_id: int = 0, role: str = "user") -> Optional[Dict]:
    """Fetch a single document if the user has access."""
    visible = get_visible_doc_ids(user_id, role)
    if visible is not None and doc_id not in visible:
        return None
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("""
                SELECT d.id, d.filename, d.file_type, d.source_name,
                       d.chunk_count, d.uploaded_by_id, d.uploaded_by_name,
                       d.version, d.parent_doc_id, d.status, d.created_at,
                       d.extra_meta, d.scope, d.scope_id
                FROM app_documents d WHERE d.id = ?
            """, doc_id)
            cols = [d[0] for d in cur.description]
            row  = cur.fetchone()
        return dict(zip(cols, row)) if row else None
    except Exception as exc:
        logger.warning("get_document failed: %s", exc)
        return None


def get_document_chunks_preview(doc_id: str, limit: int = 5) -> List[Dict]:
    """Return the first `limit` chunks of a document for preview."""
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("""
                SELECT TOP (?) chunk_index, text_content
                FROM app_document_chunks
                WHERE document_id = ?
                ORDER BY chunk_index
            """, limit, doc_id)
            return [{"chunk_index": r[0], "text": r[1]} for r in cur.fetchall()]
    except Exception as exc:
        logger.warning("get_document_chunks_preview failed: %s", exc)
        return []


def get_document_assignments(doc_id: str) -> List[Dict]:
    """Return all users a document is assigned to."""
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("""
                SELECT user_id, user_name, assigned_by_name, assigned_at
                FROM app_document_assignments WHERE doc_id = ? ORDER BY assigned_at
            """, doc_id)
            cols = [d[0] for d in cur.description]
            return [dict(zip(cols, row)) for row in cur.fetchall()]
    except Exception as exc:
        logger.warning("get_document_assignments failed: %s", exc)
        return []


def assign_document(doc_id: str, user_id: int, user_name: str,
                    assigned_by_id: int, assigned_by_name: str) -> bool:
    """
    Grant a user explicit visibility into a document via `app_document_assignments`.

    Idempotent — a MERGE that no-ops if the (doc_id, user_id) pair already exists.

    Args:
        doc_id: Document being shared.
        user_id: Target user gaining access.
        user_name: Target user's display name (denormalized for fast listing).
        assigned_by_id: ID of the user/system performing the assignment.
        assigned_by_name: Display name of the assigner.

    Returns:
        True on success, False if the SQL operation raised.
    """
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("""
                MERGE app_document_assignments AS t
                USING (SELECT ? AS doc_id, ? AS user_id) AS s
                      ON t.doc_id = s.doc_id AND t.user_id = s.user_id
                WHEN NOT MATCHED THEN
                    INSERT (doc_id, user_id, user_name, assigned_by_id, assigned_by_name)
                    VALUES (?, ?, ?, ?, ?);
            """, doc_id, user_id, doc_id, user_id, user_name, assigned_by_id, assigned_by_name)
            conn.commit()
        return True
    except Exception as exc:
        logger.warning("assign_document failed: %s", exc)
        return False


def unassign_document(doc_id: str, user_id: int) -> bool:
    """
    Revoke a user's explicit assignment to a document.

    Args:
        doc_id: Document to revoke access to.
        user_id: User whose assignment row should be removed.

    Returns:
        True on success, False if the SQL operation raised.
    """
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute(
                "DELETE FROM app_document_assignments WHERE doc_id = ? AND user_id = ?",
                doc_id, user_id)
            conn.commit()
        return True
    except Exception as exc:
        logger.warning("unassign_document failed: %s", exc)
        return False


def delete_document(document_id: str) -> Dict:
    """Remove a document — chunks from vector store + metadata + assignments + raw file + agent refs."""
    from agents.core.knowledge.vector_store import get_vector_store
    get_vector_store().delete_document(document_id)

    raw_file_path = None
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()

            # Grab raw file path from extra_meta before deleting the row
            cur.execute("SELECT extra_meta FROM app_documents WHERE id = ?", document_id)
            row = cur.fetchone()
            if row and row[0]:
                try:
                    meta = json.loads(row[0])
                    raw_file_path = meta.get("_raw_file")
                except Exception:
                    pass

            cur.execute("DELETE FROM app_document_assignments WHERE doc_id = ?", document_id)
            cur.execute("DELETE FROM app_document_chunks WHERE document_id = ?", document_id)
            cur.execute("DELETE FROM app_documents WHERE id = ?", document_id)
            conn.commit()
    except Exception as exc:
        logger.warning("delete_document SQL cleanup failed: %s", exc)

    # Delete the raw file from disk
    if raw_file_path and os.path.isfile(raw_file_path):
        try:
            os.remove(raw_file_path)
            logger.info("delete_document: removed raw file %s", raw_file_path)
        except Exception as exc:
            logger.warning("delete_document: could not remove raw file %s: %s", raw_file_path, exc)

    # Remove doc_id from any agent's search_knowledge tool config
    try:
        from app_db import get_app_db
        with get_app_db() as conn:
            cur = conn.cursor()
            cur.execute("SELECT id, tools_json FROM hub_agents")
            agents = cur.fetchall()
            for agent_id, tools_json in agents:
                if not tools_json or document_id not in tools_json:
                    continue
                try:
                    tools = json.loads(tools_json)
                    changed = False
                    for t in tools:
                        if isinstance(t, dict) and t.get("name") == "search_knowledge":
                            cfg = t.get("config") or {}
                            ids = cfg.get("document_ids", [])
                            if document_id in ids:
                                cfg["document_ids"] = [d for d in ids if d != document_id]
                                t["config"] = cfg
                                changed = True
                    if changed:
                        cur.execute(
                            "UPDATE hub_agents SET tools_json = ? WHERE id = ?",
                            json.dumps(tools), agent_id)
                except Exception:
                    pass
            conn.commit()
    except Exception as exc:
        logger.warning("delete_document: agent config cleanup failed: %s", exc)

    logger.info("delete_document: removed %s", document_id)
    return {"success": True, "document_id": document_id}
