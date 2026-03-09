"""
ppt_generator.py — Nexus · Client Presentation Generator

New in this version:
  1. Client name auto-detected from question/meeting data (not just clients.json id)
  2. Last 3 generated PPTs stored per client in clients.json as "ppt_history"
     and injected as context into every new generation
  3. Optimize endpoint: takes existing slides + user instruction → refined slides
  4. Slide 3 now web-searches competitor data live via gpt-4.1-mini
  5. bottom_message renamed to next_action — rendered as "Recommended Next Step:"
  6. Schema validation on LLM output before PPTX build
  7. Bullet text truncated to 140 chars to prevent PPTX overflow
  8. Meeting rows sorted by recency; full notes for top 5, abbreviated for next 9
  9. Email context capped at 4000 chars to prevent prompt overflow
 10. PPT history only saved when a real (non-dynamic) client is detected
"""

import os, io, json, base64, logging, re, copy, threading
from datetime import datetime
from typing import Any, Dict, List, Optional

from openai import OpenAI
from dotenv import load_dotenv
from generator_utils import make_openai_client, provider_chat
from llm_providers.models import resolve_default_model

load_dotenv()
from logging_config import get_logger
logger = get_logger(__name__)

# Email intelligence — safe fallback if module not present. logger must
# already exist by this point (see above) — the except branch below used to
# call get_logger() before it was imported, which meant this "safe fallback"
# would itself crash the whole module import if email_intelligence were ever
# actually missing, instead of degrading gracefully as intended.
try:
    from email_intelligence import (
        get_senders_for_client,
        get_emails,
        parse_date_range_from_question,
        format_emails_for_llm,
    )
    _HAS_EMAIL_INTEL = True
except ImportError:
    _HAS_EMAIL_INTEL = False
    logger.warning("email_intelligence module not found — email enrichment disabled")

# Token tracking — imported with a safe no-op fallback so ppt_generator
# works even if token_limits.py is not yet present in the project.
try:
    import token_limits as _tl
    _HAS_TOKEN_LIMITS = True
except ImportError:
    _HAS_TOKEN_LIMITS = False

from clients_db import (
    load_clients as _load_clients_db,
    get_client as _get_client_db,
    save_client as _save_client_db,
    update_ppt_history as _update_ppt_history_db,
)

OPENAI_SEARCH_MODEL = os.getenv("OPENAI_SEARCH_MODEL", "gpt-4.1-mini")


# ─────────────────────────────────────────────────────────────────────────────
# Client data helpers (SQL Server via clients_db)
# ─────────────────────────────────────────────────────────────────────────────

def _load_clients() -> List[Dict]:
    """Returns all known clients as a list of client config dicts."""
    return list(_load_clients_db().values())


def _clients_by_id() -> Dict[str, Dict]:
    """Returns all known clients keyed by client id."""
    return _load_clients_db()


def _store_ppt_history(client_id: str, slides: List[Dict],
                       question: str, is_dynamic: bool) -> None:
    """Persist last 3 PPT generations per client. Thread-safe via clients_db."""
    if is_dynamic:
        logger.info(f"⏭️  PPT history skipped — client '{client_id}' was dynamically inferred")
        return
    entry = {
        "generated_at": datetime.utcnow().isoformat(),
        "question":     question[:300],
        "slides":       slides,
    }
    _update_ppt_history_db(client_id, entry)
    logger.info(f"✅ PPT history saved for client '{client_id}'")


def _get_ppt_history(client_id: str) -> List[Dict]:
    """Returns the stored `ppt_history` list (up to 3 entries) for a client, or [] if none."""
    client = _get_client_db(client_id)
    return (client or {}).get("ppt_history", [])


# ─────────────────────────────────────────────────────────────────────────────
# Client name detection
# ─────────────────────────────────────────────────────────────────────────────

def detect_client(question: str, rows: List[Dict],
                  known_clients: List[Dict]) -> Dict:
    """
    Priority order:
      1. Exact / fuzzy match of a known client name in the question
      2. Fuzzy match in MeetingTitle fields
      3. Dynamically infer from MeetingTitle prefix (e.g. "Acme Weekly" → "Acme")
      4. Dynamically infer from the question text itself (capitalised proper nouns)
      5. Build a minimal dynamic config — never default to first known client
    """
    q_lower = question.lower()

    # 1. check known client names against question
    for c in known_clients:
        name = c.get("name", "")
        if name.lower() in q_lower:
            logger.info(f"🔍 Client matched (full name) from question: {name}")
            return c
        first_word = name.split()[0].lower()
        if len(first_word) > 3 and first_word in q_lower:
            logger.info(f"🔍 Client matched (first word) from question: {name}")
            return c
        cid = c.get("id", "")
        if cid and cid.lower() in q_lower:
            logger.info(f"🔍 Client matched (id) from question: {name}")
            return c

    # 2. check MeetingTitle in rows against known clients
    meeting_titles_raw = [
        str(r.get("MeetingTitle") or r.get("title") or "")
        for r in rows[:10]
    ]
    meeting_titles_str = " ".join(meeting_titles_raw).lower()

    for c in known_clients:
        name = c.get("name", "")
        if name.lower() in meeting_titles_str:
            logger.info(f"🔍 Client matched from meeting titles: {name}")
            return c
        first_word = name.split()[0].lower()
        if len(first_word) > 3 and first_word in meeting_titles_str:
            logger.info(f"🔍 Client matched (first word) from meetings: {name}")
            return c
        cid = c.get("id", "").lower()
        if cid and len(cid) > 1:
            title_tokens = re.findall(r'\b\w+\b', meeting_titles_str)
            if cid in title_tokens:
                logger.info(f"🔍 Client matched (id token) from meeting titles: {name}")
                return c

    # 3. Infer name dynamically from MeetingTitle prefix
    inferred_name = _infer_name_from_titles(meeting_titles_raw)

    # 4. If still nothing, try to extract a proper noun from the question
    if not inferred_name:
        inferred_name = _infer_name_from_question(question)

    if inferred_name:
        client_id = re.sub(r"[^a-z0-9]", "_", inferred_name.lower()).strip("_")
        logger.info(f"🔍 Client inferred dynamically: '{inferred_name}' (id='{client_id}')")
        return {
            "id":                      client_id,
            "name":                    inferred_name,
            "url":                     "",
            "competitors":             [],
            "our_services_for_client": [],
            "all_our_services":        [],
            "_dynamic":                True,
        }

    # 5. Absolute last resort — build an "Unknown Client" shell from the data
    logger.warning("⚠️ Could not detect or infer any client — using dynamic shell")
    return {
        "id":                      "unknown_client",
        "name":                    "Client",
        "url":                     "",
        "competitors":             [],
        "our_services_for_client": [],
        "all_our_services":        [],
        "_dynamic":                True,
    }


def _infer_name_from_titles(titles: List[str]) -> str:
    """Infers a client name from the first word of meeting titles.

    Collects the first word of each non-empty title (skipping generic
    meeting words like "weekly"/"standup"/"review" via `SKIP`). If one
    word recurs in at least 40% of titles (minimum 2), it is taken as the
    client name; otherwise falls back to the first qualifying first-word
    found in title order.

    Args:
        titles: Meeting title strings (e.g. `MeetingTitle` row values).

    Returns:
        The inferred name, or "" if no qualifying word is found.
    """
    SKIP = {
        "daily","weekly","monthly","quarterly","annual","meeting","sync",
        "call","standup","cadence","review","update","check","in","the",
        "team","project","sprint","kickoff","follow","up","status","internal",
    }
    from collections import Counter
    tokens: List[str] = []
    for t in titles:
        if not t.strip():
            continue
        first = t.strip().split()[0]
        if first.lower() not in SKIP and len(first) > 1:
            tokens.append(first)

    if not tokens:
        return ""

    most_common, count = Counter(tokens).most_common(1)[0]
    non_empty = [t for t in titles if t.strip()]
    threshold = max(2, len(non_empty) * 0.4)
    if count >= threshold:
        return most_common

    for t in titles:
        first = t.strip().split()[0] if t.strip() else ""
        if first and first.lower() not in SKIP and len(first) > 1:
            return first

    return ""


def _infer_name_from_question(question: str) -> str:
    """Infers a client name from the user's question text, as a last resort.

    Tries, in order: (1) the first quoted substring, (2) a proper-noun
    phrase following a trigger word like "for"/"about"/"client"/"regarding"/
    "on"/"of", (3) the first capitalised proper-noun phrase not entirely
    made of `STOPWORDS`.

    Args:
        question: The user's natural-language question that triggered PPT
            generation.

    Returns:
        The inferred name, or "" if nothing qualifies.
    """
    if not question:
        return ""

    STOPWORDS = {
        "show","give","me","a","an","the","ppt","presentation","slide","slides",
        "generate","create","build","make","report","for","about","client","on",
        "regarding","with","and","or","of","in","at","by","from","to","is","are",
        "what","how","why","when","where","who","status","update","progress",
        "latest","recent","new","all","my","our","this","that","these","those",
    }

    quoted = re.findall(r"['\"]([^'\"]{2,40})['\"]", question)
    if quoted:
        return quoted[0].strip()

    trigger = re.search(
        r'\b(?:for|about|client|regarding|on|of)\s+([A-Z][A-Za-z0-9&.\-]+(?:\s+[A-Z][A-Za-z0-9&.\-]+)?)',
        question
    )
    if trigger:
        candidate = trigger.group(1).strip()
        if candidate.lower() not in STOPWORDS:
            return candidate

    proper_nouns = re.findall(r'\b([A-Z][A-Za-z0-9]+(?:\s+[A-Z][A-Za-z0-9]+)*)\b', question)
    for noun in proper_nouns:
        words = noun.split()
        if all(w.lower() not in STOPWORDS for w in words) and len(noun) > 2:
            return noun

    return ""


# ─────────────────────────────────────────────────────────────────────────────
# Meeting summariser  — recency-aware, full notes for 5 most recent
# ─────────────────────────────────────────────────────────────────────────────

def _summarise_meeting_rows(rows: List[Dict]) -> str:
    """
    Sort by MeetingStart descending.
    Full notes (no truncation) for the 5 most recent meetings.
    Abbreviated (150 chars) for the next 9.
    """
    def _sort_key(r):
        return str(r.get("MeetingStart") or r.get("date") or "")

    sorted_rows = sorted(rows, key=_sort_key, reverse=True)

    lines = []
    for i, r in enumerate(sorted_rows[:14]):
        title = r.get("MeetingTitle") or r.get("title") or f"Meeting {i+1}"
        date  = str(r.get("MeetingStart") or r.get("date") or "")[:10]
        notes = r.get("NotesJson") or {}
        if isinstance(notes, str):
            try:
                notes = json.loads(notes)
            except Exception:
                notes = {}
        summary = notes.get("summary", "") if isinstance(notes, dict) else str(notes)
        summary = re.sub(r"#+ ?", "", summary)
        summary = re.sub(r"\|[^\n]+", "", summary)
        summary = summary.strip()

        # Full notes for top 5, abbreviated for the rest
        if i >= 5:
            summary = summary[:150]

        escalations = (r.get("EscalationsIssues") or "").strip()
        if i >= 5:
            escalations = escalations[:100]

        parts = []
        if escalations:
            parts.append(f"ESCALATIONS: {escalations}")
        if summary:
            parts.append(summary)
        if not parts:
            for k, v in r.items():
                if isinstance(v, str) and len(v) > 30 and k not in ("MeetingId", "ChatId"):
                    snippet = v if i < 5 else v[:200]
                    parts.append(f"{k}: {snippet}")
                    break
        body = "\n".join(parts) if parts else "(no notes)"
        lines.append(f"[{date}] {title}\n{body}")

    return "\n\n---\n\n".join(lines)


def _format_history_context(history: List[Dict]) -> str:
    """Turn stored PPT history into a readable context block."""
    if not history:
        return ""
    lines = ["PREVIOUS PRESENTATIONS FOR THIS CLIENT (newest first):"]
    for i, h in enumerate(history[:3]):
        lines.append(f"\n[PPT {i+1} — {h.get('generated_at','')[:10]}]")
        lines.append(f"Question that generated it: {h.get('question','')}")
        for s in h.get("slides", h.get("slides_summary", [])):
            col_titles = " | ".join(
                c.get("col_title", "") for c in s.get("columns", [])
            )
            lines.append(f"  • [{s.get('theme')}] {s.get('title')} → cols: {col_titles}")
    return "\n".join(lines)


# ─────────────────────────────────────────────────────────────────────────────
# Slide schema validation
# ─────────────────────────────────────────────────────────────────────────────

def _validate_slides(slides: Any) -> List[Dict]:
    """
    Ensure we have exactly 3 well-formed slide dicts.
    Pads missing bullets rather than crashing.
    Raises ValueError if the structure is unrecoverable.
    """
    if not isinstance(slides, list):
        raise ValueError(f"Expected a JSON array of slides, got {type(slides)}")
    if len(slides) != 3:
        raise ValueError(f"Expected exactly 3 slides, got {len(slides)}")

    for idx, s in enumerate(slides):
        if not isinstance(s, dict):
            raise ValueError(f"Slide {idx+1} is not an object")

        # Ensure columns exist and number exactly 3
        cols = s.get("columns")
        if not isinstance(cols, list) or len(cols) != 3:
            raise ValueError(
                f"Slide {idx+1} must have exactly 3 columns, "
                f"got {len(cols) if isinstance(cols, list) else type(cols)}"
            )

        for ci, col in enumerate(cols):
            if not isinstance(col, dict):
                raise ValueError(f"Slide {idx+1}, col {ci+1} is not an object")
            bullets = col.get("bullets", [])
            if not isinstance(bullets, list):
                col["bullets"] = ["—", "—", "—"]
            else:
                # Pad to 3 if short, truncate if somehow more
                while len(bullets) < 3:
                    bullets.append("—")
                col["bullets"] = bullets[:3]

        # Ensure next_action exists (replaces old bottom_message)
        if not s.get("next_action"):
            s["next_action"] = s.pop("bottom_message", "") or "—"

    return slides


# ─────────────────────────────────────────────────────────────────────────────
# Data quality gate
# ─────────────────────────────────────────────────────────────────────────────

_MIN_ROWS_HARD  = 1    # absolute floor — zero rows → hard reject
_MIN_ROWS_SOFT  = 3    # fewer than this → warn but still attempt if question is descriptive
_MIN_TEXT_CHARS = 40   # total text across all rows — pure-numeric datasets are flagged

def _validate_data_for_ppt(rows: List[Dict], question: str = "") -> tuple:
    """
    Returns (ok: bool, reason: str).
    ok=False means the generator should refuse and surface `reason` to the user.
    """
    n = len(rows)
    if n == 0:
        return False, "No data was returned by your query. Please ask a broader question first."
    if n < _MIN_ROWS_HARD:
        return False, f"Only {n} record found — not enough to build a 3-slide presentation."

    # Check whether rows contain any meaningful text (meeting titles, notes, names…)
    total_text = sum(
        len(str(v)) for row in rows for v in row.values()
        if isinstance(v, str) and v.strip()
    )
    if total_text < _MIN_TEXT_CHARS and n < _MIN_ROWS_SOFT:
        return False, (
            f"The {n} rows returned contain mostly numeric data with very little text. "
            "A presentation needs descriptive content — try querying meeting notes, "
            "project updates, or activity records."
        )

    if n < _MIN_ROWS_SOFT:
        return False, (
            f"Only {n} records found — the presentation would be too thin. "
            "Try broadening the date range or removing filters to return more data."
        )

    return True, ""


def _build_data_summary(rows: List[Dict], question: str = "") -> Dict:
    """Compact snapshot of query results for the optimizer context."""
    cols = list(rows[0].keys()) if rows else []
    sample = rows[:5]
    return {
        "question":  question,
        "row_count": len(rows),
        "columns":   cols,
        "sample":    sample,
    }


# ─────────────────────────────────────────────────────────────────────────────
# Bullet safety — prevent PPTX text overflow
# ─────────────────────────────────────────────────────────────────────────────

def _truncate_bullet(text: str, max_chars: int = 140) -> str:
    """Truncates text to `max_chars`, appending an ellipsis if it was cut.

    Used to keep PPTX bullet/text-box content from overflowing its shape.

    Args:
        text: Text to truncate (coerced to str).
        max_chars: Maximum length of the returned string, including the
            ellipsis character if truncation occurs.

    Returns:
        The original text if within `max_chars`, otherwise the truncated
        text with a trailing "…".
    """
    text = str(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 1] + "…"


# ─────────────────────────────────────────────────────────────────────────────
# Competitor intelligence via web search
# ─────────────────────────────────────────────────────────────────────────────

def _fetch_competitor_intelligence(
    openai_client: OpenAI,
    client_name: str,
    competitors: List[str],
    our_services: str,
) -> str:
    """
    Use gpt-4.1-mini to fetch live competitor insights.
    Returns a formatted string injected into slide 3's prompt context.
    Gracefully returns empty string on any failure.
    """
    if not competitors:
        logger.info("⏭️  No competitors listed — skipping web search for slide 3")
        return ""

    competitors_str = ", ".join(competitors)
    search_prompt = (
        f"Search the web for the latest initiatives, product launches, and strategic moves "
        f"by these companies: {competitors_str}. "
        f"Focus specifically on areas that overlap with: {our_services}. "
        f"For each competitor, summarise in 2-3 bullet points what they are doing RIGHT NOW "
        f"that is best-in-class or ahead of the market. "
        f"Be specific — cite product names, features, or announcements where possible. "
        f"Format: one section per competitor, plain text, no markdown headers."
    )

    try:
        logger.info(f"🌐 Web-searching competitor intelligence for: {competitors_str}")
        # OpenAI-only: hosted web_search_preview tool has no Bedrock equivalent,
        # so this call intentionally bypasses provider_chat()/llm_providers.
        search_response = openai_client.responses.create(
            model=OPENAI_SEARCH_MODEL,
            tools=[{"type": "web_search_preview"}],
            input=search_prompt,
        )
        # Extract text from response
        raw_text = ""
        for item in search_response.output:
            if hasattr(item, "content"):
                for block in item.content:
                    if hasattr(block, "text"):
                        raw_text += block.text
        raw_text = raw_text.strip()
        if raw_text:
            logger.info(f"✅ Competitor intelligence fetched ({len(raw_text)} chars)")
            return raw_text[:3000]  # safety cap
        logger.warning("⚠️ Web search returned empty competitor intelligence")
        return ""
    except Exception as e:
        logger.warning(f"⚠️ Competitor web search failed: {e}")
        return ""


# ─────────────────────────────────────────────────────────────────────────────
# Core Generator
# ─────────────────────────────────────────────────────────────────────────────

class PresentationGenerator:
    """Generates and refines 3-slide client-update PPTX presentations via an
    LLM (OpenAI or Amazon Bedrock, per `DEFAULT_LLM_PROVIDER` — see `llm_providers`).

    End-to-end flow: validate the BI query results are sufficient
    (`_validate_data_for_ppt`), detect/resolve the client (`detect_client`),
    enrich the prompt with meeting notes, AI insights, prior PPT history,
    optional email intelligence, and live competitor web-search data, call
    the LLM (`self.model`, default gpt-4o, via `provider_chat()`) to draft
    3 slide JSON objects via `generate()`, validate the schema
    (`_validate_slides`), persist PPT history, and render a PPTX file with
    `python-pptx` via `_build_pptx()`.
    `optimize()` re-runs a lighter version of this flow against an existing
    set of slides plus a free-text editing instruction.
    """

    def __init__(self):
        """Resolves the PPT generation model and builds a dedicated OpenAI
        client used only for the competitor web-search call (OpenAI's
        Responses API — no Bedrock equivalent). All other LLM calls in this
        class go through the provider-agnostic `provider_chat()` helper, so
        a deployment running entirely on Bedrock (no OPENAI_API_KEY at all)
        must still be able to construct this class — the competitor-search
        feature alone degrades gracefully (returns no competitor intel,
        see `_fetch_competitor_intelligence`'s own try/except) rather than
        blocking PPT generation outright.
        """
        try:
            self._openai_client = make_openai_client()
        except ValueError:
            self._openai_client = None
        self.model = resolve_default_model("OPENAI_PPT_MODEL", tier="quality")

    # ── public: generate ──────────────────────────────────────────────────────

    def generate(self, summary_points: List[str], rows: List[Dict],
                 client_config: Optional[Dict] = None,
                 question: str = "",
                 user: Optional[Dict] = None,
                 email_text: str = "") -> Dict:
        """
        Generates a fresh 3-slide client presentation from BI query results.

        summary_points: AI insight bullet strings from the chat/BI pipeline.
        rows:          Raw BI query result rows (meeting/activity records).
        client_config: optional — auto-detected from question + rows if absent.
        user:          logged-in user dict for token tracking (optional).
        email_text:    pre-formatted email context from email_intelligence.
                       Pass empty string to skip email enrichment.

        Returns:
            On success: `{"status": "success", "slides": [...], "client": {...},
            "client_id": str, "pptx_base64": str | None, "generated_at": str,
            "tokens_used": int, "email_enriched": bool, "data_summary": {...}}`.
            If `rows` fails the data-quality gate (`_validate_data_for_ppt`):
            `{"status": "insufficient_data", "message": str, "row_count": int}`.
            `pptx_base64` is None if `_build_pptx()` raised (logged, not re-raised).
        """
        # ── Data quality gate ──────────────────────────────────────────────────
        ok, reason = _validate_data_for_ppt(rows, question)
        if not ok:
            logger.warning("PPT generation rejected — insufficient data: %s", reason)
            return {"status": "insufficient_data", "message": reason, "row_count": len(rows)}
        # ─────────────────────────────────────────────────────────────────────

        known_clients = _load_clients()

        # Auto-detect client if not provided or incomplete
        if not client_config or not client_config.get("name"):
            client_config = detect_client(question, rows, known_clients)
        else:
            cid = client_config.get("id", "")
            if cid:
                full = _clients_by_id().get(cid)
                if full:
                    merged = dict(full)
                    merged.update(client_config)
                    client_config = merged

        client_id  = client_config.get("id") or \
                     re.sub(r"[^a-z0-9]", "_", client_config.get("name", "unknown").lower())
        is_dynamic = bool(client_config.get("_dynamic"))

        # Pull previous PPT history only for known clients
        history = _get_ppt_history(client_id) if not is_dynamic else []

        # Cap email context to avoid prompt overflow (~1000 tokens)
        if email_text:
            email_text = email_text[:4000]

        meeting_text = _summarise_meeting_rows(rows)
        insight_text = "\n".join(f"• {p}" for p in (summary_points or []))
        history_ctx  = _format_history_context(history)

        # Fetch live competitor intelligence for slide 3
        competitors   = client_config.get("competitors", [])
        our_services  = ", ".join(client_config.get("our_services_for_client", [])) or "data & analytics"
        competitor_intel = _fetch_competitor_intelligence(
            self._openai_client, client_config.get("name", ""), competitors, our_services
        )

        logger.info(f"🤖 Generating PPT for '{client_config.get('name')}' "
                    f"(history: {len(history)} entries, email_ctx: {bool(email_text)}, "
                    f"competitor_intel: {bool(competitor_intel)})")

        slides, prompt_used, response_used, _ppt_in, _ppt_out = self._call_openai(
            meeting_text, insight_text, client_config, history_ctx, question,
            email_text=email_text,
            competitor_intel=competitor_intel,
        )

        # Validate schema before touching PPTX builder
        try:
            slides = _validate_slides(slides)
        except ValueError as ve:
            logger.error(f"❌ Slide validation failed: {ve}")
            raise

        # Token tracking
        tokens = _record_ppt_tokens(
            user,
            call_type     = "ppt_generate",
            agent_name    = client_config.get("name", ""),
            question      = question,
            prompt        = prompt_used,
            response      = response_used,
            input_tokens  = _ppt_in,
            output_tokens = _ppt_out,
            model         = self.model,
        )
        if tokens:
            logger.info(f"[token] ppt_generate → {tokens} tokens for user {user.get('username','?')}")

        # Persist history only for real (non-dynamic) clients
        _store_ppt_history(client_id, slides, question or "", is_dynamic)

        pptx_b64 = None
        try:
            pptx_b64 = self._build_pptx(slides, client_config)
            logger.info("✅ PPTX generated")
        except Exception as e:
            logger.warning(f"⚠️ PPTX skipped: {e}")
            import traceback; traceback.print_exc()

        return {
            "status":          "success",
            "slides":          slides,
            "client":          client_config,
            "client_id":       client_id,
            "pptx_base64":     pptx_b64,
            "generated_at":    datetime.utcnow().isoformat(),
            "tokens_used":     tokens,
            "email_enriched":  bool(email_text),
            "data_summary":    _build_data_summary(rows, question),
        }

    # ── public: optimize ─────────────────────────────────────────────────────

    def optimize(self, current_slides: List[Dict], instruction: str,
                 client_config: Optional[Dict] = None,
                 user: Optional[Dict] = None,
                 data_summary: Optional[Dict] = None) -> Dict:
        """
        Refine an existing set of slides based on a user instruction.

        current_slides: The 3 slide dicts to edit (as returned by `generate()`).
        instruction:    Free-text editing instruction (validated for
                        relevance via `_validate_optimize_instruction` before
                        any LLM call is made).
        client_config:  optional client config dict (name/id/etc.).
        user:           logged-in user dict for token tracking (optional).
        data_summary: optional {question, row_count, columns, sample} from the
                      original BI query — injected into the prompt so the AI
                      knows exactly what data the presentation was built from.

        Returns:
            On success: `{"status": "success", "slides": [...], "client": {...},
            "client_id": str, "pptx_base64": str | None, "generated_at": str,
            "optimized": True, "instruction": str, "change_summary": str,
            "tokens_used": int}`.
            If the instruction is judged irrelevant/out-of-scope:
            `{"status": "invalid_instruction", "message": str}`.

        Raises:
            ValueError: If the LLM's updated slides fail `_validate_slides()`.
            json.JSONDecodeError: If the LLM response is not valid JSON.
        """
        client_config = client_config or {}
        client_name   = client_config.get("name", "the client")
        client_id     = client_config.get("id") or \
                        re.sub(r"[^a-z0-9]", "_", client_name.lower())

        logger.info(f"Optimizing PPT for '{client_name}': {instruction[:80]}")

        # ── Validate the instruction makes sense for slide editing ────────────
        validation = self._validate_optimize_instruction(
            instruction, current_slides, client_name, data_summary
        )
        if validation["status"] == "rejected":
            return {"status": "invalid_instruction", "message": validation["reason"]}
        # ─────────────────────────────────────────────────────────────────────

        slides_json = json.dumps(current_slides, indent=2)

        # Build data context block so AI knows what the slides were built from
        data_ctx = ""
        if data_summary:
            orig_q   = data_summary.get("question", "")
            row_cnt  = data_summary.get("row_count", 0)
            cols     = ", ".join(data_summary.get("columns", []))
            sample   = json.dumps(data_summary.get("sample", [])[:3], default=str)
            data_ctx = (
                f"\nDATA CONTEXT (what this presentation was built from):\n"
                f"  Original query: \"{orig_q}\"\n"
                f"  Rows fetched: {row_cnt}\n"
                f"  Columns: {cols}\n"
                f"  Sample rows: {sample}\n"
                "Use this context to make edits more specific and data-grounded.\n"
            )

        prompt = f"""You are a senior consultant at Nexus editing a 3-slide client presentation for {client_name}.
{data_ctx}
INSTRUCTION FROM USER: {instruction}

CURRENT SLIDES (JSON):
{slides_json}

Rules:
- Keep the same JSON structure (theme, title, subtitle, context_label, context_icon,
  context_message, columns[col_title, col_subtitle, bullets×3], next_action)
- Apply the instruction precisely — change only what was asked
- Leave slides not mentioned by the instruction unchanged
- If the instruction names a specific slide (e.g. "slide 2", "competitive slide"), update only that one
- If the instruction is general or says "all slides", update all three
- Ground edits in the actual data context above — be specific, not generic
- Do NOT add or remove slides
- next_action must be one concrete action Nexus should propose to {client_name} this week

Return a JSON object with exactly two keys:
  "slides"  – the updated array of 3 slide objects
  "summary" – 1-2 sentences describing what you changed and why (plain text, no markdown)

No markdown fences. Raw JSON only."""

        result = provider_chat(
            system="You are a senior management consultant. Always respond with raw JSON only — no markdown, no explanation.",
            user_prompt=prompt,
            model=self.model,
            temperature=0.3,
            max_tokens=3600,
        )
        raw = result.content.strip()

        # Parse: expect {"slides": [...], "summary": "..."}
        clean = re.sub(r"```(?:json)?", "", raw).replace("```", "").strip()
        parsed = json.loads(clean)
        if isinstance(parsed, dict):
            slides        = parsed.get("slides", [])
            change_summary = parsed.get("summary", "")
        else:
            # Fallback: AI returned a bare array
            slides        = parsed
            change_summary = ""

        # Validate schema
        try:
            slides = _validate_slides(slides)
        except ValueError as ve:
            logger.error(f"Slide validation failed (optimize): {ve}")
            raise

        logger.info(f"Optimized {len(slides)} slides via OpenAI ({self.model})")

        _opt_in    = result.usage.prompt_tokens
        _opt_out   = result.usage.completion_tokens
        tokens = _record_ppt_tokens(
            user,
            call_type     = "ppt_optimize",
            agent_name    = client_name,
            question      = instruction,
            prompt        = prompt,
            response      = raw,
            input_tokens  = _opt_in,
            output_tokens = _opt_out,
            model         = self.model,
        )

        is_dynamic = bool(client_config.get("_dynamic"))
        _store_ppt_history(client_id, slides, f"[OPTIMIZED] {instruction}", is_dynamic)

        pptx_b64 = None
        try:
            pptx_b64 = self._build_pptx(slides, client_config)
        except Exception as e:
            logger.warning(f"PPTX skipped: {e}")

        return {
            "status":         "success",
            "slides":         slides,
            "client":         client_config,
            "client_id":      client_id,
            "pptx_base64":    pptx_b64,
            "generated_at":   datetime.utcnow().isoformat(),
            "optimized":      True,
            "instruction":    instruction,
            "change_summary": change_summary,
            "tokens_used":    tokens,
        }

    def _validate_optimize_instruction(
        self,
        instruction: str,
        current_slides: List[Dict],
        client_name: str,
        data_summary: Optional[Dict],
    ) -> Dict:
        """
        Ask the LLM if the instruction is a valid slide-editing request.
        Returns {"status": "ok"} or {"status": "rejected", "reason": "..."}.
        """
        slide_titles = " | ".join(s.get("title", "") for s in current_slides)
        orig_q = (data_summary or {}).get("question", "")
        data_line = f'The presentation was built from the query: "{orig_q}"' if orig_q else ""

        prompt = (
            f'A user is editing a 3-slide client presentation for "{client_name}".\n'
            f'Slide titles: {slide_titles}\n'
            f'{data_line}\n\n'
            f'User instruction: "{instruction}"\n\n'
            'Is this a valid slide-editing instruction? Classify as one of:\n'
            '- valid        : changes tone, content, emphasis, focus, length, style, or slide-specific edits\n'
            '- irrelevant   : completely unrelated to the presentation (e.g. weather, jokes, cooking)\n'
            '- out_of_scope : asks for something impossible in this format (add charts, embed video, etc.)\n\n'
            'Reply with JSON only: {"verdict": "valid|irrelevant|out_of_scope", "reason": "one sentence"}'
        )
        try:
            result = provider_chat(
                system="",
                user_prompt=prompt,
                model=resolve_default_model("OPENAI_PPT_VALIDATE_MODEL", tier="fast"),
                temperature=0,
                max_tokens=80,
            )
            raw  = result.content.strip()
            raw  = re.sub(r"```(?:json)?", "", raw).replace("```", "").strip()
            data = json.loads(raw)
            verdict = data.get("verdict", "valid")
            reason  = data.get("reason", "")
            if verdict in ("irrelevant", "out_of_scope"):
                friendly = {
                    "irrelevant":   f"That instruction doesn't seem related to editing the presentation. Please describe a change you'd like to make to the slides (e.g. \"make slide 1 more concise\", \"update the next action\", \"focus on ROI\").",
                    "out_of_scope": f"That's not something that can be done in this slide format. {reason}",
                }
                return {"status": "rejected", "reason": friendly[verdict]}
        except Exception as e:
            logger.warning("Instruction validation failed (%s) — proceeding anyway", e)
        return {"status": "ok"}

    # ── OpenAI: generate prompt ───────────────────────────────────────────────

    def _call_openai(self, meeting_text: str, insight_text: str,
                     client_config: Dict, history_ctx: str,
                     question: str, email_text: str = "",
                     competitor_intel: str = "") -> tuple:
        """Builds the slide-generation prompt and calls `self.model` (default gpt-4o).

        Assembles a single prompt covering meeting notes, AI insights,
        engagement/competitor context, optional email intelligence and
        prior-PPT history, then asks the model for exactly 3 slide JSON
        objects (themes: "delivery", "process", "competitive"). When
        `competitor_intel` is non-empty (from `_fetch_competitor_intelligence`),
        slide 3's prompt section uses that live web data instead of the
        model's own background knowledge.

        Args:
            meeting_text: Recency-sorted meeting summary from `_summarise_meeting_rows()`.
            insight_text: Bullet-formatted AI insight strings.
            client_config: Client config dict (name/url/competitors/services).
            history_ctx: Formatted prior-PPT context from `_format_history_context()`.
            question: The user's original question, used to focus the narrative.
            email_text: Optional formatted email-intelligence context.
            competitor_intel: Optional live competitor web-search text from
                `_fetch_competitor_intelligence()`.

        Returns:
            `(slides, prompt, raw_response, input_tokens, output_tokens)` —
            `slides` is the parsed JSON (expected to be a list of 3 slide
            dicts; not yet schema-validated), `prompt`/`raw_response` are
            the exact strings sent/received for token-recording and
            debugging, and the token counts come from the provider's usage object.
        """

        client_name  = client_config.get("name", "the client")
        client_url   = client_config.get("url", "")
        competitors  = ", ".join(client_config.get("competitors", [])) or "key market players"
        our_services = ", ".join(client_config.get("our_services_for_client", [])) or "data & analytics"
        all_services = ", ".join(client_config.get("all_our_services", [])) or ""

        history_section = (
            f"\n\n{history_ctx}\n\n"
            "IMPORTANT: Use the previous presentations above to avoid repeating the same "
            "content. Advance the narrative — show progress, new findings, or deeper "
            "recommendations compared to what was already presented.\n"
        ) if history_ctx else ""

        question_section = (
            f"\nUSER QUESTION THAT TRIGGERED THIS: {question}\n"
            "Tailor the presentation angle and emphasis to what the user specifically asked about.\n"
        ) if question else ""

        email_section = (
            f"\n\nEMAIL INTELLIGENCE (use this to add communication-level specifics):\n"
            f"{email_text}\n"
            "When email data is available, reference specific email subjects, discussions, "
            "and action items in the slide bullets. This makes the presentation far more "
            "credible and specific.\n"
        ) if email_text else ""

        # Slide 3 competitor section — live web data if available, fallback to general prompt
        if competitor_intel:
            slide3_context = f"""LIVE COMPETITOR INTELLIGENCE (fetched from the web — use this directly):
{competitor_intel}

For slide 3, use the above real competitor data to:
- Show exactly what {competitors} are doing right now that is best-in-class
- Frame each column around one competitor or one capability theme
- For each item, show what that competitor is doing AND what Nexus can offer {client_name} in response
- Make it forward-looking and opportunity-framed, not a deficit analysis"""
        else:
            slide3_context = (
                f"For slide 3, draw on your knowledge of {competitors} to identify "
                f"best-in-class practices relevant to {our_services} that {client_name} could adopt."
            )

        prompt = f"""You are a senior consultant at Nexus preparing a 3-slide client update for {client_name} ({client_url}).
{question_section}
MEETING DATA:
{meeting_text}

AI INSIGHTS:
{insight_text}

CURRENT ENGAGEMENT: {our_services}
ADDITIONAL SERVICES: {all_services}
COMPETITORS: {competitors}
{email_section}{history_section}
Generate exactly 3 slide objects as a JSON array. Be SPECIFIC — reference real names, systems, tasks from the meeting data. No generic content.

SLIDE 1 theme="delivery" — title should be "{client_name}: Current improvements"
  Focus on concrete delivery progress from the meetings.
  context_label = "What changed"
  context_icon = "↑"

SLIDE 2 theme="process" — title should be "{client_name}: New work that can add value"
  Focus on 3 specific next-step opportunities emerging from the current project.
  context_label = "What we can unlock next"
  context_icon = "+"

SLIDE 3 theme="competitive" — title should be "{client_name}: Where the market is heading"
  {slide3_context}
  context_label = "The opportunity"
  context_icon = "→"

Each slide:
- subtitle: one sharp line framing the slide
- context_message: one sentence on the arc
- columns: exactly 3, each with col_title (2-4 words), col_subtitle (one sharp line), bullets (exactly 3 strings of 1-2 sentences each, grounded in the data)
- next_action: one specific action Nexus should propose to {client_name} this week (NOT prefixed with "Suggested client message:")

Return ONLY a valid JSON array of exactly 3 objects. No markdown fences.
[
  {{
    "theme": "delivery",
    "title": "...", "subtitle": "...",
    "context_label": "...", "context_icon": "...", "context_message": "...",
    "columns": [
      {{"col_title":"...","col_subtitle":"...","bullets":["...","...","..."]}},
      {{"col_title":"...","col_subtitle":"...","bullets":["...","...","..."]}},
      {{"col_title":"...","col_subtitle":"...","bullets":["...","...","..."]}}
    ],
    "next_action": "..."
  }},
  ...
]"""

        result = provider_chat(
            system="You are a senior management consultant. Always respond with raw JSON only — no markdown, no explanation.",
            user_prompt=prompt,
            model=self.model,
            temperature=0.3,
            max_tokens=3500,
        )
        raw      = result.content.strip()
        slides   = _safe_json(raw)
        _in_tok  = result.usage.prompt_tokens
        _out_tok = result.usage.completion_tokens
        logger.info(f"✅ Got {len(slides)} slides from OpenAI ({self.model})")
        return slides, prompt, raw, _in_tok, _out_tok

    # ── PPTX builder ─────────────────────────────────────────────────────────

    def _build_pptx(self, slides: List[Dict],
                    client_config: Dict) -> Optional[str]:
        """Renders validated slide dicts into a 16:9 PPTX deck using python-pptx.

        Each slide gets a theme-colored layout (one of "delivery"/"process"/
        "competitive" from `THEMES`, defaulting to "delivery") with a title,
        subtitle, divider, a two-part context banner (icon pill + message),
        3 bullet-point columns, and a "Recommended Next Step" footer built
        from `next_action`. Text is truncated via `_truncate_bullet()` to
        avoid PPTX shape overflow.

        Args:
            slides: List of slide dicts already validated by
                `_validate_slides()` (theme/title/subtitle/context_*/
                columns[col_title, col_subtitle, bullets]/next_action).
            client_config: Client config dict (currently unused inside this
                method, but kept in the signature for caller symmetry with
                `generate()`/`optimize()`).

        Returns:
            The rendered PPTX file, base64-encoded as a UTF-8 string —
            this is the `pptx_base64` value returned by `generate()` and
            `optimize()`, and the value `scheduler_service.py` decodes
            before attaching it to an email.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        THEMES = {
            "delivery": {
                "primary":    RGBColor(0x60, 0x80, 0xAE),
                "light":      RGBColor(0xEB, 0xF2, 0xF9),
                "mid":        RGBColor(0x48, 0x67, 0x95),
                "bg":         RGBColor(0xF5, 0xF9, 0xFD),
                "context_bg": RGBColor(0xEB, 0xF2, 0xF9),
            },
            "process": {
                "primary":    RGBColor(0xAD, 0x81, 0x38),
                "light":      RGBColor(0xF8, 0xF3, 0xEB),
                "mid":        RGBColor(0x48, 0x67, 0x95),
                "bg":         RGBColor(0xF5, 0xF9, 0xFD),
                "context_bg": RGBColor(0xF8, 0xF3, 0xEB),
            },
            "competitive": {
                "primary":    RGBColor(0x58, 0x92, 0x74),
                "light":      RGBColor(0xEC, 0xF5, 0xEF),
                "mid":        RGBColor(0x48, 0x67, 0x95),
                "bg":         RGBColor(0xF5, 0xF9, 0xFD),
                "context_bg": RGBColor(0xEC, 0xF5, 0xEF),
            },
        }

        WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
        DARK      = RGBColor(0x37, 0x48, 0x5E)
        MUTED     = RGBColor(0x6F, 0x7F, 0x94)
        DIVIDER   = RGBColor(0xDD, 0xE6, 0xF1)
        FOOTER_BG = RGBColor(0xF8, 0xFA, 0xFD)

        def I(v): return Inches(v)  # noqa: E743 - shorthand for Inches() used throughout this builder

        prs = Presentation()
        prs.slide_width  = I(13.333)
        prs.slide_height = I(7.5)
        blank = prs.slide_layouts[6]

        def rect(slide, l, t, w, h, fill_rgb, radius=False, line_rgb=None, line_pt=0):
            """Adds a filled rectangle (or rounded rectangle if radius=True) shape to `slide`."""
            sp = slide.shapes.add_shape(5 if radius else 1, I(l), I(t), I(w), I(h))
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill_rgb
            if line_rgb:
                sp.line.color.rgb = line_rgb
                sp.line.width = Pt(line_pt)
            else:
                sp.line.fill.background()
            return sp

        def oval(slide, l, t, w, h, fill_rgb):
            """Adds a filled oval shape to `slide` (used for numbered column bullets/icons)."""
            sp = slide.shapes.add_shape(9, I(l), I(t), I(w), I(h))
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill_rgb
            sp.line.fill.background()
            return sp

        def txt(slide, text, l, t, w, h,
                size=12, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, italic=False):
            """Adds a single-paragraph text box to `slide`, truncating text to 200 chars."""
            # Cap text length to prevent overflow
            text = _truncate_bullet(str(text), max_chars=200)
            txb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
            tf  = txb.text_frame
            tf.word_wrap = True
            tf.margin_left   = I(0.05)
            tf.margin_top    = I(0.03)
            tf.margin_right  = I(0.03)
            tf.margin_bottom = I(0.02)
            p   = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text        = text
            run.font.size   = Pt(size)
            run.font.bold   = bold
            run.font.italic = italic
            if color:
                run.font.color.rgb = color
            return txb

        def txt_multi(slide, lines, l, t, w, h, size=12, color=DARK):
            """Adds a multi-paragraph text box to `slide`, one paragraph per line (each truncated to 140 chars)."""
            txb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
            tf  = txb.text_frame
            tf.word_wrap = True
            tf.margin_left   = I(0.05)
            tf.margin_top    = I(0.03)
            tf.margin_right  = I(0.04)
            tf.margin_bottom = I(0.02)
            for i, line in enumerate(lines):
                # Truncate each bullet to prevent card overflow
                safe_line = _truncate_bullet(str(line), max_chars=140)
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(3)
                run = p.add_run()
                run.text           = safe_line
                run.font.size      = Pt(size)
                run.font.color.rgb = color
            return txb

        for slide_data in slides:
            theme_key = slide_data.get("theme", "delivery")
            C         = THEMES.get(theme_key, THEMES["delivery"])
            primary   = C["primary"]
            light     = C["light"]
            mid       = C["mid"]
            bg        = C["bg"]
            ctx_bg    = C["context_bg"]

            sl = prs.slides.add_slide(blank)
            rect(sl, 0, 0, 13.333, 7.5, WHITE)

            # Title
            txt(sl, slide_data.get("title", ""),
                0.55, 0.22, 11.0, 0.52, size=28, bold=True, color=primary)
            # Subtitle
            txt(sl, slide_data.get("subtitle", ""),
                0.58, 0.72, 11.0, 0.28, size=11, color=MUTED)
            # Divider
            rect(sl, 0.55, 1.06, 10.90, 0.025, DIVIDER)

            # Context banner — left pill
            rect(sl, 0.72, 1.32, 2.25, 0.64, ctx_bg, radius=True)
            oval(sl, 0.86, 1.46, 0.34, 0.34, primary)
            txt(sl, slide_data.get("context_icon", "✦"),
                0.88, 1.50, 0.28, 0.22, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
            txt(sl, slide_data.get("context_label", ""),
                1.24, 1.48, 1.55, 0.22, size=12, bold=True, color=primary)
            # Context banner — right message
            rect(sl, 3.15, 1.32, 7.50, 0.64, bg, radius=True)
            txt(sl, slide_data.get("context_message", ""),
                3.35, 1.42, 7.10, 0.42, size=12, bold=True, color=DARK)

            # Three columns
            columns     = slide_data.get("columns", [{}, {}, {}])
            col_xs      = [0.72, 4.12, 7.52]
            col_fills   = [light, bg, light]
            col_headers = [primary, mid, primary]

            CARD_T = 2.18; CARD_W = 3.28; CARD_H = 3.72
            HDR_H  = 0.50; OVAL_W = 0.28

            for ci, (cx, col_data) in enumerate(zip(col_xs, columns)):
                hdr_color = col_headers[ci]
                rect(sl, cx, CARD_T, CARD_W, CARD_H, col_fills[ci], radius=True)
                rect(sl, cx + 0.03, CARD_T + 0.03, CARD_W - 0.06, HDR_H,
                     hdr_color, radius=True)
                ov_x = cx + 0.16; ov_y = CARD_T + 0.12
                oval(sl, ov_x, ov_y, OVAL_W, OVAL_W, WHITE)
                txt(sl, str(ci + 1), ov_x, ov_y + 0.12, OVAL_W, 0.20,
                    size=10, bold=True, color=hdr_color, align=PP_ALIGN.CENTER)
                txt(sl, col_data.get("col_title", ""),
                    cx + 0.52, CARD_T + 0.14, CARD_W - 0.62, 0.26,
                    size=14, bold=True, color=WHITE)
                txt(sl, col_data.get("col_subtitle", ""),
                    cx + 0.18, CARD_T + HDR_H + 0.16, CARD_W - 0.26, 0.40,
                    size=10, bold=True, color=hdr_color)
                txt_multi(sl, col_data.get("bullets", []),
                          cx + 0.18, CARD_T + HDR_H + 0.62,
                          CARD_W - 0.26, CARD_H - HDR_H - 0.72,
                          size=11, color=DARK)

            # Footer — "Recommended Next Step:" label + next_action text
            next_action = slide_data.get("next_action", "")
            footer_text = f"Recommended Next Step: {next_action}" if next_action else ""
            rect(sl, 1.00, 6.35, 9.95, 0.46, FOOTER_BG, radius=True)
            txt(sl, footer_text,
                1.15, 6.44, 9.65, 0.32, size=11, bold=True, color=DARK)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return base64.b64encode(buf.read()).decode("utf-8")


# ─────────────────────────────────────────────────────────────────────────────
# Token recording helper
# ─────────────────────────────────────────────────────────────────────────────

def _record_ppt_tokens(user: Optional[Dict], call_type: str,
                       agent_name: str, question: str,
                       prompt: str, response: str,
                       input_tokens: int = 0, output_tokens: int = 0,
                       model: str = "") -> int:
    """Records OpenAI token usage for a PPT generate/optimize call against the requesting user.

    No-op (returns 0) if the `token_limits` module is unavailable, `user`
    is falsy, or `user` has no "id" (e.g. background/admin calls).

    Args:
        user: current_user dict (must contain "id") to attribute usage to.
        call_type: One of "ppt_generate" / "ppt_optimize" (used as the
            token-usage category).
        agent_name: Client name, used as the recorded "agent" label.
        question: The triggering question/instruction, for audit context.
        prompt: Full prompt text sent to the LLM.
        response: Raw LLM response text.
        input_tokens: Prompt token count from the OpenAI usage object, if known.
        output_tokens: Completion token count from the OpenAI usage object, if known.

    Returns:
        The token count recorded, or 0 if recording was skipped/failed.
    """
    if not _HAS_TOKEN_LIMITS or not user or not user.get("id"):
        return 0
    try:
        _in  = int(input_tokens  or 0)
        _out = int(output_tokens or 0)
        _tot = _in + _out or None
        return _tl.record_tokens(
            user,
            call_type=call_type,
            agent_name=agent_name,
            question=question,
            prompt=prompt,
            response=response,
            actual_tokens=_tot,
            input_tokens=_in,
            output_tokens=_out,
            model=model,
        )
    except Exception as e:
        logger.warning(f"PPT token tracking failed ({call_type}): {e}")
        return 0


# ─────────────────────────────────────────────────────────────────────────────
# Safe JSON parser
# ─────────────────────────────────────────────────────────────────────────────

def _safe_json(raw: str) -> Any:
    """Parses an LLM JSON response, stripping markdown fences and unwrapping a "slides" key.

    Args:
        raw: Raw LLM response text, possibly wrapped in ```json ... ``` fences.

    Returns:
        `parsed["slides"]` if the parsed JSON is a dict containing a
        "slides" key, otherwise the parsed JSON as-is (expected to already
        be a bare list of slide dicts).

    Raises:
        json.JSONDecodeError: If `raw` (after fence stripping) is not valid JSON.
    """
    clean = re.sub(r"```(?:json)?", "", raw).replace("```", "").strip()
    parsed = json.loads(clean)
    if isinstance(parsed, dict) and "slides" in parsed:
        return parsed["slides"]
    return parsed